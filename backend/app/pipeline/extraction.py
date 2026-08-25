"""Text extraction for supported document formats.

Each extractor is pure (bytes -> text) so it is trivially unit-testable and free
of framework/storage concerns.
"""

from __future__ import annotations

import csv
import io

from app.core.logging import get_logger

logger = get_logger("engineergpt.pipeline.extract")


def extract_text(extension: str, data: bytes) -> str:
    """Dispatch to the right extractor by file extension."""
    handler = {
        ".txt": _from_txt,
        ".csv": _from_csv,
        ".pdf": _from_pdf,
        ".docx": _from_docx,
        ".pptx": _from_pptx,
        ".xlsx": _from_xlsx,
    }.get(extension)
    if handler is None:
        logger.info("extract_unsupported", extra={"extra_fields": {"ext": extension}})
        return ""
    try:
        return handler(data)
    except Exception:
        logger.exception("extract_failed", extra={"extra_fields": {"ext": extension}})
        return ""


def _from_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _from_csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(", ".join(row) for row in reader)


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _from_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.extend(p.text for p in shape.text_frame.paragraphs)
    return "\n".join(lines)


def _from_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            lines.append(", ".join("" if c is None else str(c) for c in row))
    return "\n".join(lines)
